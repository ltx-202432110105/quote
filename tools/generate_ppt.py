#!/usr/bin/env python3
from __future__ import annotations

import argparse
import glob
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


FONT_NAME = "Microsoft YaHei"
MONO_FONT = "Consolas"

COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PRIMARY = RGBColor(0x2F, 0x6F, 0xED)
COLOR_TEXT = RGBColor(0x11, 0x18, 0x27)
COLOR_SUBTEXT = RGBColor(0x4B, 0x55, 0x63)
COLOR_LINE = RGBColor(0xAA, 0xB4, 0xC3)
COLOR_CODE_BG = RGBColor(0x1F, 0x29, 0x37)
COLOR_CODE_BORDER = RGBColor(0x2F, 0x3C, 0x52)
COLOR_GLASS = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GLASS_BORDER = RGBColor(0xD5, 0xDF, 0xEA)
COLOR_DECOR = RGBColor(0x8B, 0xA3, 0xC7)

TEAM_NAME = "我们叫什么名字"
SCHOOL_NAME = "浙江师范大学"
TRACK_NAME = "A类"

MAX_BULLETS_PER_SLIDE = 5
MAX_CODE_LINES = 8
MAX_TABLE_ROWS = 6
TARGET_TOTAL_SLIDES = 20
TARGET_CONTENT_SLIDES = TARGET_TOTAL_SLIDES - 1
MAX_AGENDA_ITEMS = 8
AGENDA_HIGHLIGHT_COUNT = 3
TITLE_LEN_NORMAL = 18
SUMMARY_TITLE = "总结与展望"
MERGE_TABLE_PENALTY = 4
MERGE_CODE_PENALTY = 4
PLACEHOLDER_BULLET = "待补充内容"
GO_KEYWORDS = ("go", "golang", "gin", "grpc", "编译", "go语言")
TECH_KEYWORD_ICON_MAP = {
    "go": "go.png",
    "golang": "go.png",
    "go语言": "go.png",
    "gin": "go.png",
    "grpc": "go.png",
    "go-zero": "go.png",
    "docker": "docker.png",
    "mysql": "mysql.png",
    "redis": "redis.png",
    "k8s": "kubernetes.png",
    "kubernetes": "kubernetes.png",
    "postgres": "postgresql.png",
    "postgresql": "postgresql.png",
    "nginx": "nginx.png",
    "github": "github.png",
}
ICON_DISPLAY_NAMES = {
    "go": "GO",
    "docker": "DOCKER",
    "mysql": "MYSQL",
    "redis": "REDIS",
    "kubernetes": "K8S",
    "postgresql": "POSTGRESQL",
    "nginx": "NGINX",
    "github": "GITHUB",
}
DEFAULT_PANEL_TAGS = ["GO + GIN", "MYSQL + REDIS"]
COVER_TECH_ORDER = ["GO", "DOCKER", "MYSQL", "REDIS", "K8S", "POSTGRESQL", "NGINX", "GITHUB"]
SUMMARY_BULLETS = ["实现职业规划核心链路闭环", "强化推荐质量与可解释反馈", "下一步推进多模型协同优化"]
QA_BULLETS = ["期待专家指导与建议", "交流：技术方案与落地实践", f"队伍：{TEAM_NAME}"]
PROCESS_KEYWORDS = ("流程", "步骤", "阶段", "演进", "架构")
CHAPTER_KEYWORDS = ("总结", "感谢", "目录", "展望", "第", "功能")
MARKDOWN_LANGUAGE_MARKERS = {"text", "yaml", "json", "sql", "go", "bash"}
TECH_LAYOUT_KEYWORDS = ("架构", "技术", "部署", "实现")
REFERENCE_PDF = "比赛用的最终ppt.pdf"
BACKGROUND_GLOB = "assets/backgrounds/*"
ICONS_DIR = Path("assets/icons")
COVER_TECH_COUNT = 4
BACKGROUND_LAYOUT_RIGHT = "right"
BACKGROUND_LAYOUT_FULL = "full"
BACKGROUND_LAYOUT_DIAGONAL = "diagonal"
BACKGROUND_LAYOUTS = (BACKGROUND_LAYOUT_RIGHT, BACKGROUND_LAYOUT_FULL, BACKGROUND_LAYOUT_DIAGONAL)
GO_LOGO_X = Inches(10.7)
GO_LOGO_Y = Inches(1.52)
GO_LOGO_WIDTH = Inches(2.0)
LAYOUT_DEFAULT = "A"
LAYOUT_DOUBLE = "B"
LAYOUT_CHAPTER = "C"
LAYOUT_TIMELINE = "D"
LAYOUT_AGENDA = "E"
TEXTURE_TRANSPARENCY_EMPHASIS = 0.89
TEXTURE_TRANSPARENCY_DEFAULT = 0.94
TEXTURE_DOT_TRANSPARENCY = 0.85
TEXTURE_ARC_TRANSPARENCY = 0.86
# Dot coordinate tuples are (x, y) in inches relative to slide origin.
TEXTURE_DOTS_EMPHASIS = [(8.55, 1.7), (9.1, 2.2), (10.35, 3.0), (11.2, 4.25)]
TEXTURE_DOTS_DEFAULT = [(10.9, 3.15), (11.35, 3.7)]


@dataclass
class TableData:
    header_left: str
    header_right: str
    rows: list[tuple[str, str]] = field(default_factory=list)


@dataclass
class Section:
    title: str
    bullets: list[str] = field(default_factory=list)
    tables: list[TableData] = field(default_factory=list)
    code_blocks: list[list[str]] = field(default_factory=list)


@dataclass
class SlideUnit:
    title: str
    bullets: list[str] = field(default_factory=list)
    table: TableData | None = None
    code_block: list[str] | None = None
    is_transition: bool = False
    layout_hint: str | None = None


def set_para_style(paragraph, size: int, color: RGBColor, bold: bool = False, align: PP_ALIGN | None = None) -> None:
    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    if align is not None:
        paragraph.alignment = align


def parse_markdown_sections(path: Path) -> list[Section]:
    text = path.read_text(encoding="utf-8")
    parts = re.split(r"^##\s+", text, flags=re.MULTILINE)
    sections: list[Section] = []

    for part in parts:
        part = part.strip()
        if not part:
            continue
        lines = part.splitlines()
        section = Section(title=lines[0].strip())

        i = 1
        while i < len(lines):
            line = lines[i].rstrip()
            stripped = line.strip()

            if not stripped or stripped == "---" or stripped.lower() in MARKDOWN_LANGUAGE_MARKERS:
                i += 1
                continue

            if stripped.startswith("```"):
                code: list[str] = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    code.append(lines[i].rstrip())
                    i += 1
                if code:
                    section.code_blocks.append(code)
                i += 1
                continue

            if stripped.startswith("|") and stripped.endswith("|"):
                table_lines: list[str] = []
                while i < len(lines):
                    row = lines[i].strip()
                    if row.startswith("|") and row.endswith("|"):
                        table_lines.append(row)
                        i += 1
                    else:
                        break
                table = parse_markdown_table(table_lines)
                if table:
                    section.tables.append(table)
                continue

            clean = cleanup_text(stripped)
            if clean:
                section.bullets.append(clean)
            i += 1

        deduped: list[str] = []
        seen = set()
        for bullet in section.bullets:
            if bullet in seen:
                continue
            seen.add(bullet)
            deduped.append(bullet)
        section.bullets = deduped
        sections.append(section)

    return sections


def parse_markdown_table(lines: list[str]) -> TableData | None:
    rows: list[list[str]] = []
    for line in lines:
        cols = [c.strip() for c in line.strip("|").split("|")]
        if cols and all(set(c) <= {"-", ":"} for c in cols):
            continue
        rows.append(cols)

    if len(rows) < 2:
        return None

    head = rows[0]
    body = rows[1:]
    left = head[0] if len(head) > 0 else "项目"
    right = head[1] if len(head) > 1 else "说明"

    data = TableData(header_left=left, header_right=right)
    for row in body:
        l = row[0] if len(row) > 0 else ""
        r = row[1] if len(row) > 1 else ""
        if l or r:
            data.rows.append((cleanup_text(l), cleanup_text(r)))
    return data


def cleanup_text(text: str) -> str:
    text = re.sub(r"^[-*]\s+", "", text)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"<[^>]+>", "", text)
    return text.strip()


def chunk_bullets(bullets: list[str], chunk_size: int) -> list[list[str]]:
    return [bullets[i:i + chunk_size] for i in range(0, len(bullets), chunk_size)]


def normalize_and_shorten_title(title: str, max_len: int = TITLE_LEN_NORMAL) -> str:
    """Normalize title text (remove page/pagination markers, collapse spaces) then truncate."""
    t = cleanup_text(title).replace("：", " ").replace("—", " ")
    t = re.sub(r"^第\s*\d+\s*页\s*[-—:：]?\s*", "", t)
    t = re.sub(r"（\d+/\d+）", "", t)
    t = re.sub(r"\s+", " ", t).strip()
    if len(t) <= max_len:
        return t
    return f"{t[:max_len].rstrip()}…"


def summarize_bullets(bullets: list[str], max_items: int = MAX_BULLETS_PER_SLIDE) -> list[str]:
    """Limit bullet count to max_items and append an overflow summary note when needed."""
    cleaned: list[str] = []
    for bullet in bullets:
        item = cleanup_text(bullet)
        if item:
            cleaned.append(item)
    if not cleaned:
        return [PLACEHOLDER_BULLET]
    if len(cleaned) <= max_items:
        return cleaned
    summarized = cleaned[:max_items - 1]
    overflow_count = len(cleaned) - (max_items - 1)
    summarized.append(f"其余 {overflow_count} 项已整合到现场讲解")
    return summarized


def split_section(section: Section) -> list[SlideUnit]:
    units: list[SlideUnit] = []
    title = normalize_and_shorten_title(section.title, max_len=22)

    if section.bullets:
        chunks = chunk_bullets(section.bullets, MAX_BULLETS_PER_SLIDE)
        total = len(chunks)
        for idx, chunk in enumerate(chunks, start=1):
            if len(section.bullets) > MAX_BULLETS_PER_SLIDE and total > 1:
                chunk_title = f"{title}（{idx}/{total}）"
            else:
                chunk_title = title
            units.append(SlideUnit(title=chunk_title, bullets=summarize_bullets(chunk)))

    if not section.bullets and not section.tables and not section.code_blocks:
        units.append(SlideUnit(title=title, bullets=[PLACEHOLDER_BULLET]))

    for table in section.tables[:1]:
        units.append(SlideUnit(title=title, table=table))

    for code in section.code_blocks[:1]:
        units.append(SlideUnit(title=title, code_block=code[:MAX_CODE_LINES]))

    if should_use_transition(section.title, section.bullets) and units:
        units[0].is_transition = True

    return units


def should_use_transition(title: str, bullets: list[str]) -> bool:
    txt = f"{title} {' '.join(bullets[:2])}"
    return any(k in txt for k in CHAPTER_KEYWORDS) and len(bullets) <= 2


def select_layout(unit: SlideUnit) -> str:
    if unit.layout_hint:
        return unit.layout_hint
    title_text = unit.title.lower()
    all_text = f"{unit.title} {' '.join(unit.bullets)}".lower()
    if unit.is_transition:
        return LAYOUT_CHAPTER
    if any(k in all_text for k in PROCESS_KEYWORDS) and unit.bullets:
        return LAYOUT_TIMELINE
    if unit.table or len(unit.bullets) >= 5:
        return LAYOUT_DOUBLE
    if any(k in title_text for k in TECH_LAYOUT_KEYWORDS):
        return LAYOUT_DEFAULT
    return LAYOUT_DEFAULT


def need_go_logo(unit: SlideUnit) -> bool:
    txt = f"{unit.title} {' '.join(unit.bullets)}".lower()
    return any(k in txt for k in GO_KEYWORDS)


def detect_tech_icons(unit: SlideUnit, icon_dir: Path = ICONS_DIR) -> list[Path]:
    text = f"{unit.title} {' '.join(unit.bullets)}".lower()
    hits: list[Path] = []
    for kw, icon_name in TECH_KEYWORD_ICON_MAP.items():
        if kw in text:
            p = icon_dir / icon_name
            if p.exists() and p not in hits:
                hits.append(p)
    return hits[:4]


def extract_tech_keywords(unit: SlideUnit) -> list[str]:
    tags: list[str] = []
    for icon_path in detect_tech_icons(unit):
        label = ICON_DISPLAY_NAMES.get(icon_path.stem, icon_path.stem.upper())
        if label not in tags:
            tags.append(label)
    return tags[:2]


def resolve_backgrounds() -> list[Path]:
    files = [Path(p) for p in sorted(glob.glob(BACKGROUND_GLOB))]
    return [p for p in files if p.is_file() and p.suffix.lower() in {".jpg", ".jpeg", ".png"}]


def pick_background(backgrounds: list[Path], index: int) -> Path | None:
    if not backgrounds:
        return None
    return backgrounds[index % len(backgrounds)]


def pick_background_layout(index: int) -> str:
    return BACKGROUND_LAYOUTS[index % len(BACKGROUND_LAYOUTS)]


def add_background_hero(slide, prs: Presentation, background_path: Path | None, layout: str) -> None:
    if not background_path or not background_path.exists():
        return

    if layout == BACKGROUND_LAYOUT_RIGHT:
        slide.shapes.add_picture(str(background_path), Inches(8.55), Inches(0), width=Inches(4.78), height=prs.slide_height)
        veil = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.42), Inches(0), Inches(4.95), prs.slide_height)
        veil.fill.solid()
        veil.fill.fore_color.rgb = COLOR_WHITE
        veil.fill.transparency = 0.22
        veil.line.fill.background()
    elif layout == BACKGROUND_LAYOUT_DIAGONAL:
        slide.shapes.add_picture(str(background_path), Inches(5.2), Inches(0), width=Inches(8.15), height=prs.slide_height)
        cut = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.85), Inches(0.05), Inches(2.65), Inches(7.35))
        cut.fill.solid()
        cut.fill.fore_color.rgb = RGBColor(0xE6, 0xEE, 0xFA)
        cut.fill.transparency = 0.12
        cut.line.fill.background()
        veil = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(0), Inches(8.25), prs.slide_height)
        veil.fill.solid()
        veil.fill.fore_color.rgb = COLOR_WHITE
        veil.fill.transparency = 0.26
        veil.line.fill.background()
    else:
        slide.shapes.add_picture(str(background_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        veil = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        veil.fill.solid()
        veil.fill.fore_color.rgb = COLOR_WHITE
        veil.fill.transparency = 0.28
        veil.line.fill.background()


def apply_theme(
    slide,
    prs: Presentation,
    background_path: Path | None = None,
    background_layout: str = BACKGROUND_LAYOUT_RIGHT,
    emphasize_texture: bool = False,
) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG

    add_background_hero(slide, prs, background_path, background_layout)

    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    glow.fill.solid()
    glow.fill.fore_color.rgb = COLOR_WHITE
    glow.fill.transparency = 0.72
    glow.line.fill.background()

    light = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(0.2), Inches(6.4), Inches(7.0))
    light.fill.solid()
    light.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
    light.fill.transparency = 0.9
    light.line.fill.background()

    top_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06))
    top_strip.fill.solid()
    top_strip.fill.fore_color.rgb = COLOR_PRIMARY
    top_strip.line.fill.background()

    add_background_texture(slide, emphasize=emphasize_texture)
    add_hud_decor(slide)


def add_hud_decor(slide) -> None:
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(0.35), Inches(0.02), Inches(0.36))
    left.fill.solid()
    left.fill.fore_color.rgb = COLOR_PRIMARY
    left.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(0.35), Inches(0.46), Inches(0.02))
    top.fill.solid()
    top.fill.fore_color.rgb = COLOR_PRIMARY
    top.line.fill.background()

    right_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.55), Inches(0.35), Inches(0.46), Inches(0.02))
    right_top.fill.solid()
    right_top.fill.fore_color.rgb = COLOR_LINE
    right_top.line.fill.background()

    right_side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.99), Inches(0.35), Inches(0.02), Inches(0.36))
    right_side.fill.solid()
    right_side.fill.fore_color.rgb = COLOR_LINE
    right_side.line.fill.background()


def add_background_texture(slide, emphasize: bool = False) -> None:
    transparency = TEXTURE_TRANSPARENCY_EMPHASIS if emphasize else TEXTURE_TRANSPARENCY_DEFAULT
    h_lines = [Inches(1.9), Inches(3.25)] if emphasize else [Inches(4.75)]
    v_lines = [Inches(8.35), Inches(9.85)] if emphasize else [Inches(11.25)]
    dots = TEXTURE_DOTS_EMPHASIS if emphasize else TEXTURE_DOTS_DEFAULT

    for y in h_lines:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, Inches(6.3), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_DECOR
        line.fill.transparency = transparency
        line.line.fill.background()

    for x in v_lines:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.55), Inches(0.01), Inches(3.7))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_DECOR
        line.fill.transparency = transparency
        line.line.fill.background()

    for x, y in dots:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.05), Inches(0.05))
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_DECOR
        dot.fill.transparency = TEXTURE_DOT_TRANSPARENCY
        dot.line.fill.background()

    circuit_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.1), Inches(2.15), Inches(1.6), Inches(0.01))
    circuit_h.fill.solid()
    circuit_h.fill.fore_color.rgb = COLOR_DECOR
    circuit_h.fill.transparency = transparency
    circuit_h.line.fill.background()
    circuit_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.68), Inches(2.15), Inches(0.01), Inches(1.2))
    circuit_v.fill.solid()
    circuit_v.fill.fore_color.rgb = COLOR_DECOR
    circuit_v.fill.transparency = transparency
    circuit_v.line.fill.background()
    circuit_h2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.68), Inches(3.34), Inches(1.2), Inches(0.01))
    circuit_h2.fill.solid()
    circuit_h2.fill.fore_color.rgb = COLOR_DECOR
    circuit_h2.fill.transparency = transparency
    circuit_h2.line.fill.background()

    for idx in range(4):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.25), Inches(1.55 + idx * 0.36), Inches(0.5), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_DECOR
        bar.fill.transparency = 0.9 if idx % 2 == 0 else 0.84
        bar.line.fill.background()

    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(10.55), Inches(4.45), Inches(1.65), Inches(1.1))
    arc.fill.background()
    arc.line.color.rgb = COLOR_DECOR
    arc.line.width = Pt(1)
    arc.line.transparency = TEXTURE_ARC_TRANSPARENCY


def add_title_block(slide, title: str) -> None:
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(8.8), Inches(0.7))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    set_para_style(p, 28, COLOR_TEXT, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.17), Inches(3.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()


def add_footer(slide, prs: Presentation, page_no: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), prs.slide_height - Inches(0.56), prs.slide_width - Inches(1.4), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LINE
    line.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.75), prs.slide_height - Inches(0.42), Inches(8.5), Inches(0.24))
    tf = left.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"队伍：{TEAM_NAME}  ·  学校：{SCHOOL_NAME}  ·  赛道：{TRACK_NAME}"
    set_para_style(p, 10, COLOR_SUBTEXT)

    right = slide.shapes.add_textbox(prs.slide_width - Inches(1.0), prs.slide_height - Inches(0.42), Inches(0.4), Inches(0.24))
    rtf = right.text_frame
    rtf.clear()
    p2 = rtf.paragraphs[0]
    p2.text = str(page_no)
    set_para_style(p2, 10, COLOR_SUBTEXT, align=PP_ALIGN.RIGHT)


def add_glass_card(
    slide,
    x,
    y,
    w,
    h,
    fill_color: RGBColor = COLOR_GLASS,
    border_color: RGBColor = COLOR_GLASS_BORDER,
    transparency: float = 0.22,
):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(1), y + Pt(1), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0xE9, 0xEE, 0xF4)
    shadow.fill.transparency = 0.65
    shadow.line.fill.background()
    shadow.adjustments[0] = 0.14

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.fill.transparency = transparency
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    card.adjustments[0] = 0.14
    return card


def add_bullets(slide, bullets: list[str], x, y, w, h, size=16) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        set_para_style(p, size, COLOR_TEXT if idx < 2 else COLOR_SUBTEXT)
        p.space_after = Pt(8)
        p.line_spacing = 1.25


def render_table_card(slide, table: TableData, x, y, w, h) -> None:
    add_glass_card(slide, x, y, w, h, transparency=0.18)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.12), y + Inches(0.10), w - Inches(0.24), Inches(0.42))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0xEC, 0xF2, 0xFF)
    header.line.fill.background()

    left_w = (w - Inches(0.24)) * 0.32
    h1 = slide.shapes.add_textbox(x + Inches(0.16), y + Inches(0.14), left_w, Inches(0.3))
    h1tf = h1.text_frame
    h1tf.clear()
    p = h1tf.paragraphs[0]
    p.text = table.header_left
    set_para_style(p, 13, COLOR_PRIMARY, bold=True)

    h2 = slide.shapes.add_textbox(x + Inches(0.16) + left_w + Inches(0.12), y + Inches(0.14), w - Inches(0.5) - left_w, Inches(0.3))
    h2tf = h2.text_frame
    h2tf.clear()
    p2 = h2tf.paragraphs[0]
    p2.text = table.header_right
    set_para_style(p2, 13, COLOR_PRIMARY, bold=True)

    row_start = y + Inches(0.58)
    row_h = min(Inches(0.62), (h - Inches(0.75)) / max(1, len(table.rows[:MAX_TABLE_ROWS])))
    for idx, (left_text, right_text) in enumerate(table.rows[:MAX_TABLE_ROWS]):
        ry = row_start + idx * row_h
        if idx > 0:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.16), ry - Inches(0.04), w - Inches(0.32), Inches(0.01))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0xE3, 0xE8, 0xEF)
            sep.line.fill.background()

        lbox = slide.shapes.add_textbox(x + Inches(0.16), ry, left_w, row_h)
        ltf = lbox.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lp.text = left_text
        set_para_style(lp, 12, COLOR_TEXT, bold=True)

        rbox = slide.shapes.add_textbox(x + Inches(0.16) + left_w + Inches(0.12), ry, w - Inches(0.5) - left_w, row_h)
        rtf = rbox.text_frame
        rtf.clear()
        rp = rtf.paragraphs[0]
        rp.text = right_text
        set_para_style(rp, 12, COLOR_SUBTEXT)


def render_code_card(slide, code_lines: list[str], x, y, w, h) -> None:
    card = add_glass_card(slide, x, y, w, h, fill_color=COLOR_CODE_BG, border_color=COLOR_CODE_BORDER, transparency=0)
    # Tune corner radius ratio for a cleaner digital panel style.
    card.adjustments[0] = 0.12

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.14), y + Inches(0.10), Inches(1.0), Inches(0.28))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(0x34, 0x42, 0x59)
    tag.line.fill.background()
    ttf = tag.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = "CODE"
    set_para_style(tp, 10, RGBColor(0xD4, 0xDE, 0xEF), bold=True, align=PP_ALIGN.CENTER)

    box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.46), w - Inches(0.36), h - Inches(0.56))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for idx, line in enumerate(code_lines[:MAX_CODE_LINES]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        if not p.runs:
            p.add_run()
        run = p.runs[0]
        run.font.name = MONO_FONT
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xE7, 0xEF, 0xFF)
        p.space_after = Pt(2)


def add_right_tech_panel(slide, unit: SlideUnit, panel_label: str = "TECH PANEL") -> None:
    panel = add_glass_card(slide, Inches(9.65), Inches(1.45), Inches(2.85), Inches(5.35), fill_color=RGBColor(0xF8, 0xFB, 0xFF), transparency=0.26)
    ptf = panel.text_frame
    ptf.clear()
    p = ptf.paragraphs[0]
    p.text = panel_label
    set_para_style(p, 11, COLOR_SUBTEXT, align=PP_ALIGN.CENTER)

    module = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.05), Inches(2.55), Inches(2.05), Inches(2.05))
    module.fill.background()
    module.line.color.rgb = COLOR_DECOR
    module.line.transparency = 0.46
    module.line.width = Pt(1.2)
    module.adjustments[0] = 0.08

    for y in [2.9, 3.25, 3.6, 3.95]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.25), Inches(y), Inches(1.65), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_DECOR
        line.fill.transparency = 0.8
        line.line.fill.background()

    for x in [10.45, 11.05, 11.65]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.75), Inches(0.01), Inches(1.55))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_DECOR
        line.fill.transparency = 0.84
        line.line.fill.background()

    tech_tags = extract_tech_keywords(unit) or DEFAULT_PANEL_TAGS
    for idx, txt in enumerate(tech_tags[:2]):
        t = slide.shapes.add_textbox(Inches(9.95), Inches(4.8 + idx * 0.34), Inches(2.3), Inches(0.2))
        tf = t.text_frame
        tf.clear()
        p2 = tf.paragraphs[0]
        p2.text = txt
        set_para_style(p2, 10, COLOR_SUBTEXT)
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.8), Inches(4.88 + idx * 0.34), Inches(1.25), Inches(0.01))
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLOR_DECOR
        marker.fill.transparency = 0.78
        marker.line.fill.background()

    glow = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(9.92), Inches(2.9), Inches(2.25), Inches(1.75))
    glow.fill.background()
    glow.line.color.rgb = COLOR_PRIMARY
    glow.line.transparency = 0.86
    glow.line.width = Pt(1)

    icons = detect_tech_icons(unit)
    for idx, icon in enumerate(icons[:4]):
        x = Inches(9.95 + (idx % 2) * 1.25)
        y = Inches(1.95 + (idx // 2) * 1.05)
        slide.shapes.add_picture(str(icon), x, y, width=Inches(0.8), height=Inches(0.8))


def render_layout_a(slide, unit: SlideUnit) -> None:
    add_glass_card(slide, Inches(0.8), Inches(1.45), Inches(8.6), Inches(5.35), transparency=0.2)
    add_bullets(slide, unit.bullets, Inches(1.05), Inches(1.72), Inches(8.1), Inches(4.8), size=16)
    add_right_tech_panel(slide, unit, panel_label="TECH STACK")


def render_layout_b(slide, unit: SlideUnit) -> None:
    add_glass_card(slide, Inches(0.8), Inches(1.45), Inches(4.25), Inches(5.35), transparency=0.2)
    add_glass_card(slide, Inches(5.15), Inches(1.45), Inches(4.25), Inches(5.35), transparency=0.2)
    add_right_tech_panel(slide, unit, panel_label="TECH CARD")

    mid = (len(unit.bullets) + 1) // 2
    left = unit.bullets[:mid]
    right = unit.bullets[mid:]

    add_bullets(slide, left, Inches(1.05), Inches(1.72), Inches(3.9), Inches(4.9), size=14)
    add_bullets(slide, right, Inches(5.4), Inches(1.72), Inches(3.9), Inches(4.9), size=14)


def render_layout_e(slide, unit: SlideUnit) -> None:
    add_glass_card(slide, Inches(0.9), Inches(1.45), Inches(8.45), Inches(5.3), transparency=0.2)
    add_right_tech_panel(slide, unit, panel_label="AGENDA")

    for idx, item in enumerate(unit.bullets[:MAX_AGENDA_ITEMS], start=1):
        y = Inches(1.8 + (idx - 1) * 0.58)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(0.48), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0xE9, 0xF1, 0xFF)
        badge.line.fill.background()
        btf = badge.text_frame
        btf.clear()
        bp = btf.paragraphs[0]
        bp.text = str(idx)
        set_para_style(bp, 10, COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.78), y + Inches(0.16), Inches(0.32), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_DECOR
        line.fill.transparency = 0.78
        line.line.fill.background()

        box = slide.shapes.add_textbox(Inches(2.2), y - Inches(0.03), Inches(6.7), Inches(0.4))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = item
        set_para_style(p, 15, COLOR_TEXT if idx <= AGENDA_HIGHLIGHT_COUNT else COLOR_SUBTEXT)


def render_layout_c(slide, unit: SlideUnit) -> None:
    badge_text = "CHAPTER"
    page_match = re.search(r"第\s*(\d+)\s*(页|章)", unit.title)
    if page_match:
        badge_text = f"PAGE {page_match.group(1)}"

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.78), Inches(2.05), Inches(0.56))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0xE9, 0xF0, 0xFF)
    badge.fill.transparency = 0.18
    badge.line.fill.background()
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = badge_text
    set_para_style(p, 14, COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    glass = add_glass_card(slide, Inches(0.95), Inches(2.45), Inches(8.45), Inches(2.55), transparency=0.22)
    glass.line.color.rgb = RGBColor(0xD7, 0xE2, 0xEF)
    add_right_tech_panel(slide, unit, panel_label="SECTION")

    tbox = slide.shapes.add_textbox(Inches(1.25), Inches(2.8), Inches(7.8), Inches(2.0))
    ttf = tbox.text_frame
    ttf.clear()
    p2 = ttf.paragraphs[0]
    p2.text = unit.title
    set_para_style(p2, 40, COLOR_TEXT, bold=True)

    if unit.bullets:
        sbox = slide.shapes.add_textbox(Inches(1.25), Inches(5.28), Inches(8.8), Inches(0.9))
        stf = sbox.text_frame
        stf.clear()
        p3 = stf.paragraphs[0]
        p3.text = unit.bullets[0]
        set_para_style(p3, 18, COLOR_SUBTEXT)

    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(10.8), Inches(4.9), Inches(1.7), Inches(1.2))
    arc.fill.background()
    arc.line.color.rgb = COLOR_PRIMARY
    arc.line.transparency = 0.86
    arc.line.width = Pt(1.2)


def render_layout_d(slide, unit: SlideUnit) -> None:
    add_glass_card(slide, Inches(0.8), Inches(1.55), Inches(8.6), Inches(5.15), transparency=0.2)
    add_right_tech_panel(slide, unit, panel_label="FLOW")

    steps = unit.bullets[:5]
    if not steps:
        steps = ["阶段一", "阶段二", "阶段三"]

    total = len(steps)
    for idx, step in enumerate(steps):
        if total == 1:
            cx = Inches(4.75)
        else:
            cx = Inches(1.05 + idx * (6.9 / (total - 1)))
        cy = Inches(3.15)

        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.52), Inches(0.52))
        node.fill.solid()
        node.fill.fore_color.rgb = COLOR_PRIMARY
        node.line.fill.background()
        ntf = node.text_frame
        ntf.clear()
        np = ntf.paragraphs[0]
        np.text = str(idx + 1)
        set_para_style(np, 12, COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

        tbox = slide.shapes.add_textbox(cx - Inches(0.38), cy + Inches(0.65), Inches(1.35), Inches(1.2))
        ttf = tbox.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        tp.text = step
        set_para_style(tp, 12, COLOR_SUBTEXT, align=PP_ALIGN.CENTER)

        if idx < total - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, cx + Inches(0.62), cy + Inches(0.18), Inches(1.0), Inches(0.17))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0xB8, 0xC5, 0xDD)
            arrow.line.fill.background()


def add_cover(
    prs: Presentation,
    title: str,
    subtitle: str,
    logo_path: Path,
    background_path: Path | None = None,
    background_layout: str = BACKGROUND_LAYOUT_FULL,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_theme(slide, prs, background_path=background_path, background_layout=background_layout, emphasize_texture=True)

    hero = add_glass_card(slide, Inches(0.95), Inches(1.2), Inches(11.4), Inches(5.25), fill_color=COLOR_WHITE, transparency=0.16)
    hero.line.color.rgb = RGBColor(0xD7, 0xE0, 0xEF)

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.25), Inches(1.55), Inches(2.0), Inches(0.52))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(0xE8, 0xF1, 0xFF)
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "A类赛道路演"
    set_para_style(p, 13, COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    title_box = slide.shapes.add_textbox(Inches(1.25), Inches(2.2), Inches(8.2), Inches(1.9))
    ttf = title_box.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = title
    set_para_style(tp, 38, COLOR_TEXT, bold=True)

    sub = slide.shapes.add_textbox(Inches(1.25), Inches(4.2), Inches(8.6), Inches(1.2))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    set_para_style(sp, 18, COLOR_SUBTEXT)

    info = add_glass_card(slide, Inches(9.0), Inches(2.0), Inches(2.95), Inches(2.8), fill_color=RGBColor(0xFB, 0xFC, 0xFF), transparency=0.2)
    itf = info.text_frame
    itf.clear()
    p1 = itf.paragraphs[0]
    p1.text = f"队伍\n{TEAM_NAME}"
    set_para_style(p1, 12, COLOR_TEXT, bold=True)
    p2 = itf.add_paragraph()
    p2.text = f"学校\n{SCHOOL_NAME}"
    set_para_style(p2, 12, COLOR_SUBTEXT)
    p3 = itf.add_paragraph()
    p3.text = f"赛道\n{TRACK_NAME}"
    set_para_style(p3, 12, COLOR_SUBTEXT)

    cover_tech = " / ".join(COVER_TECH_ORDER[:COVER_TECH_COUNT])
    cover_unit = SlideUnit(title=title, bullets=[subtitle, cover_tech])
    add_right_tech_panel(slide, cover_unit, panel_label="COVER / TECH")

    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(9.4), Inches(4.35), width=Inches(1.9))


def extract_cover(sections: list[Section]) -> tuple[str, str]:
    if not sections:
        return "比赛项目汇报", "简约科技数码风演示文稿"
    title = "比赛项目汇报"
    subtitle = "简约科技数码风演示文稿"
    for line in sections[0].bullets:
        if "项目名称" in line and "：" in line:
            title = line.split("：", 1)[1].strip() or title
        if "一句话" in line and "：" in line:
            subtitle = line.split("：", 1)[1].strip() or subtitle
    return title, subtitle


def normalize_unit(unit: SlideUnit) -> SlideUnit:
    """Create a compact unit with shortened title and bounded bullets for slide rendering."""
    return SlideUnit(
        title=normalize_and_shorten_title(unit.title, max_len=TITLE_LEN_NORMAL),
        bullets=summarize_bullets(unit.bullets),
        table=unit.table,
        code_block=unit.code_block,
        is_transition=unit.is_transition,
        layout_hint=unit.layout_hint,
    )


def merge_units_for_limit(units: list[SlideUnit], target: int) -> list[SlideUnit]:
    """Merge adjacent low-complexity units until the unit count reaches target."""
    merged = list(units)
    while len(merged) > target:
        best_idx = 0
        best_score = float("inf")
        for i in range(len(merged) - 1):
            score = len(merged[i].bullets) + len(merged[i + 1].bullets)
            if merged[i].table or merged[i + 1].table:
                score += MERGE_TABLE_PENALTY
            if merged[i].code_block or merged[i + 1].code_block:
                score += MERGE_CODE_PENALTY
            if score < best_score:
                best_score = score
                best_idx = i
        a = merged[best_idx]
        b = merged[best_idx + 1]
        title = normalize_and_shorten_title(a.title, max_len=20)
        bullets = summarize_bullets([*a.bullets, *b.bullets], MAX_BULLETS_PER_SLIDE)
        merged_unit = SlideUnit(
            title=title,
            bullets=bullets,
            table=a.table or b.table,
            code_block=a.code_block or b.code_block,
            is_transition=a.is_transition or b.is_transition,
            layout_hint=a.layout_hint or b.layout_hint,
        )
        merged = merged[:best_idx] + [merged_unit] + merged[best_idx + 2:]
    return merged


def build_agenda_unit(units: list[SlideUnit]) -> SlideUnit:
    """Build a fixed agenda slide from the first key unit titles."""
    agenda_items = [normalize_and_shorten_title(u.title, max_len=20) for u in units[:MAX_AGENDA_ITEMS]]
    if len(units) > MAX_AGENDA_ITEMS:
        agenda_items[-1] = SUMMARY_TITLE
    return SlideUnit(title="目录总览", bullets=agenda_items or ["项目概览", "系统设计", "技术实现", "总结展望"], layout_hint=LAYOUT_AGENDA)


def build_padding_units() -> list[SlideUnit]:
    """Return fallback design slides used to pad content to the target page count."""
    # These fallback templates are defaults for this competition deck and should be customized per project.
    return [
        SlideUnit(title="关键亮点与创新", bullets=["（可按项目替换）智能评估链路闭环", "（可按项目替换）岗位匹配准确率提升", "（可按项目替换）全链路可解释反馈", "（可按项目替换）多模型融合增强鲁棒性"], layout_hint=LAYOUT_DEFAULT),
        SlideUnit(title="系统架构总览", bullets=["（可按项目替换）接入层与网关", "（可按项目替换）核心服务引擎", "（可按项目替换）数据层与缓存层", "（可按项目替换）监控与治理体系"], layout_hint=LAYOUT_TIMELINE),
        SlideUnit(title="技术栈设计", bullets=["（可按项目替换）后端与接口框架", "（可按项目替换）数据库与缓存组件", "（可按项目替换）容器化与编排能力", "（可按项目替换）协作与交付平台"], layout_hint=LAYOUT_DOUBLE),
        SlideUnit(title="部署与运维", bullets=["（可按项目替换）发布与灰度策略", "（可按项目替换）日志指标链路追踪", "（可按项目替换）弹性扩缩容治理", "（可按项目替换）安全与备份方案"], layout_hint=LAYOUT_DEFAULT),
        SlideUnit(title="演示流程", bullets=["（可按项目替换）入口与鉴权", "（可按项目替换）核心处理流程", "（可按项目替换）结果反馈流程", "（可按项目替换）展示与收尾"], layout_hint=LAYOUT_TIMELINE),
    ]


def build_content_units(sections: list[Section]) -> list[SlideUnit]:
    """Build fixed 19 content slides: agenda + 16 body + summary + Q&A."""
    raw_units: list[SlideUnit] = []
    for sec in sections:
        raw_units.extend(split_section(sec))
    raw_units = [normalize_unit(u) for u in raw_units]
    body_target = 16
    raw_units = merge_units_for_limit(raw_units, body_target)

    pads = build_padding_units()
    pad_idx = 0
    while len(raw_units) < body_target:
        # Cycle through fallback templates using modulo when multiple padding pages are required.
        raw_units.append(pads[pad_idx % len(pads)])
        pad_idx += 1

    body_units = merge_units_for_limit(raw_units, body_target)[:body_target]
    content: list[SlideUnit] = [build_agenda_unit(body_units)]
    content.extend(body_units)
    content.append(SlideUnit(title=SUMMARY_TITLE, bullets=SUMMARY_BULLETS, layout_hint=LAYOUT_CHAPTER, is_transition=True))
    content.append(SlideUnit(title="Q&A / 感谢", bullets=QA_BULLETS, layout_hint=LAYOUT_CHAPTER, is_transition=True))
    return content[:TARGET_CONTENT_SLIDES]


def build_ppt(md_paths: list[Path], output_path: Path, logo_path: Path) -> None:
    sections: list[Section] = []
    for md in md_paths:
        sections.extend(parse_markdown_sections(md))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    backgrounds = resolve_backgrounds()

    title, subtitle = extract_cover(sections)
    cover_bg = pick_background(backgrounds, 0)
    add_cover(prs, title, subtitle, logo_path, background_path=cover_bg, background_layout=BACKGROUND_LAYOUT_FULL)

    content_units = build_content_units(sections)
    page_no = 1
    for idx, unit in enumerate(content_units, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        layout = select_layout(unit)
        bg = pick_background(backgrounds, idx)
        bg_layout = pick_background_layout(idx)
        apply_theme(slide, prs, background_path=bg, background_layout=bg_layout, emphasize_texture=(layout == LAYOUT_CHAPTER))
        if layout not in {LAYOUT_CHAPTER, LAYOUT_AGENDA}:
            add_title_block(slide, unit.title)
        if unit.table:
            render_table_card(slide, unit.table, Inches(0.95), Inches(1.55), Inches(8.45), Inches(5.1))
            add_right_tech_panel(slide, unit, panel_label="DATA PANEL")
        elif unit.code_block:
            render_code_card(slide, unit.code_block, Inches(0.95), Inches(1.55), Inches(8.45), Inches(5.1))
            add_right_tech_panel(slide, unit, panel_label="CODE PANEL")
        elif layout == LAYOUT_DOUBLE:
            render_layout_b(slide, unit)
        elif layout == LAYOUT_CHAPTER:
            render_layout_c(slide, unit)
        elif layout == LAYOUT_TIMELINE:
            render_layout_d(slide, unit)
        elif layout == LAYOUT_AGENDA:
            render_layout_e(slide, unit)
        else:
            render_layout_a(slide, unit)

        add_footer(slide, prs, page_no)
        page_no += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def resolve_markdowns(patterns: list[str]) -> list[Path]:
    files: list[Path] = []
    for pattern in patterns:
        for matched in glob.glob(pattern):
            p = Path(matched)
            if p.is_file() and p.suffix.lower() == ".md":
                files.append(p)
    files = sorted(set(files))
    if not files:
        raise FileNotFoundError(f"No markdown files found matching: {patterns}")
    return files


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate competition PPTX with gray-tech style")
    parser.add_argument("--inputs", nargs="+", default=["CareerPlanner*.md", "docs/**/*.md"])
    parser.add_argument("--output", default="slides/competition.pptx")
    parser.add_argument("--go-logo", default="assets/icons/go.png")
    args = parser.parse_args()

    md_paths = resolve_markdowns(args.inputs)
    output = Path(args.output)
    go_logo = Path(args.go_logo)

    build_ppt(md_paths, output, go_logo)

    print("Markdown sources:")
    for md in md_paths:
        print(f"- {md}")
    print(f"Reference PDF: {REFERENCE_PDF}")
    print(f"Go logo: {go_logo}")
    print(f"Output PPTX: {output}")


if __name__ == "__main__":
    main()
